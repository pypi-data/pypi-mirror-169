#!/usr/bin/python
# -*- coding: utf-8 -*-
# editor: mufei(ypdh@qq.com tel:+086 15712150708)
'''
Mufei _ __ ___   ___   ___  / _| ___(_)
| '_ ` _ \ / _ \ / _ \| |_ / _ \ |
| | | | | | (_) | (_) |  _|  __/ |
|_| |_| |_|\___/ \___/|_|  \___|_|
'''

import warnings

try:
    from pptx import Presentation 
except ImportError:
    warnings.warn('Please Import python-pptx TO pip')
except NameError:
    warnings.warn('Please Import python-pptx TO pip')

try:
    from .cplugin_fix import _StringIO, StringIO, cPlugin_Fix
except (ImportError,ValueError):
    from cplugin_fix import _StringIO, StringIO, cPlugin_Fix   
from moofei._find import _py, _strtypes,_find_func,_search_func



class Plugin_Fix_Pptx(cPlugin_Fix):
             
    def all_texts(self):
        Ls = []
        prs = Presentation(self.fp)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:#判断Shape是否含有文本框
                    text_frame = shape.text_frame
                    #replace_text(text_frame)#调用replace_text函数实现文本替换
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            #run.text = 
                            Ls.append(run.text)
                                
                if shape.has_table:#判断Shape是否含有表格
                    table = shape.table
                    for cell in table.iter_cells():#遍历表格的cell
                        text_frame = cell.text_frame
                        #replace_text(text_frame))#调用replace_text函数实现文本替换
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                #run.text = 
                                Ls.append(run.text)
        #prs.save(filename_save)#保存

        text = '\n'.join(Ls)
        return text
        
        
