import tempfile
from collections import Counter
import matplotlib.pyplot as plt
from slugify import slugify
from PIL import Image
from scriptine import path
import pptx
from pptx.util import Inches, Cm
from xdat import xsettings


DEFAULT_THEME = xsettings.XDAT_ROOT.joinpath('media', 'default_theme.pptx')
assert DEFAULT_THEME.exists(), DEFAULT_THEME


class Presentation:
    def __init__(self, theme=DEFAULT_THEME):
        self.theme = theme
        self.prs = pptx.Presentation(self.theme)

    def add_slide(self, layout_name, **kwargs):
        add_slide(self.prs, layout_name, **kwargs)

    def print_layout(self):
        print_layout(self.theme)

    def save(self, out_path):
        self.prs.save(out_path)

    @classmethod
    def capture_image(cls):
        return Img()


class Img:
    """
    Everything's in inches...
    """
    DPI = 80

    def __init__(self, tight_layout=True):
        assert xsettings.CACHE_PATH is not None, "must set xsettings.CACHE_PATH"
        tmp_folder = xsettings.CACHE_PATH.joinpath('xpptx')
        tmp_folder.ensure_dir()

        self.img_path = tempfile.NamedTemporaryFile(suffix='.png', delete=False, dir=tmp_folder).name

        if tight_layout:
            plt.tight_layout()

        plt.savefig(self.img_path, pad_inches=0)
        plt.close('all')
        img = Image.open(self.img_path)
        w, h = img.size
        self.width = w/self.DPI
        self.height = h/self.DPI

    def box(self, width, height):
        rw = self.width / width
        rh = self.height / height
        rmax = max(rh, rw)
        return self.width/rmax, self.height/rmax

    def __str__(self):
        return self.img_path

    def __repr__(self):
        return self.img_path

    def __lt__(self, other):
        return False


def _slug(text):
    return slugify(text, separator="_")


def _slug_dict(input_dict):
    counts = Counter()
    out_dict = dict()
    for text, v in input_dict.items():
        try:
            parts = text.split()
            int(parts[-1])
            parts = parts[:-1]
            text = " ".join(parts)
        except:
            pass

        try:
            parts = text.split()
            if parts[-1].lower() == 'placeholder':
                parts = parts[:-1]
            text = " ".join(parts)
        except:
            pass

        text = _slug(text)
        counts[text] += 1

        if counts[text] > 1:
            text = f"{text}_{counts[text]}"

        out_dict[text] = v

    return out_dict


def _get_layout(prs):
    layouts = _slug_dict({l.name: l for l in prs.slide_layouts})
    return layouts


def _get_placeholders(slide):
    placeholders = _slug_dict({p.name: p for p in slide.placeholders})
    return placeholders


def add_slide(prs, layout_name, **kwargs):
    layouts = _get_layout(prs)
    assert layout_name in layouts, f"{layout_name} not in: {sorted(layouts)}"

    layout = layouts[layout_name]
    slide = prs.slides.add_slide(layout)

    placeholders = _get_placeholders(slide)

    for k, v in kwargs.items():
        assert k in placeholders, f"{k} not in: {sorted(placeholders)}"
        p = placeholders[k]

        if isinstance(v, str):
            p.text = v

        elif isinstance(v, Img):
            w,h = v.box(p.width, p.height)
            w = int(w)
            h = int(h)
            slide.shapes.add_picture(str(v), p.left, p.top, height=h)
            p.text = ' '

        else:
            raise TypeError(type(v))

    for k in set(placeholders) - set(kwargs):
        p = placeholders[k]
        p.text = ' '

    return


def print_layout(template_path=None):
    prs = pptx.Presentation(template_path)
    layouts = _get_layout(prs)
    for layout_name in sorted(layouts):
        print(f"- {layout_name}")
        layout = layouts[layout_name]
        slide = prs.slides.add_slide(layout)
        placeholders = _get_placeholders(slide)
        for pl in sorted(placeholders):
            print(f"  + {pl}")

    return


if __name__ == "__main__":
    xsettings.PROJECT_NAME = 'xdat'
    xsettings.updated_config()
    plt.scatter([1, 2], [3, 4])
    i = Img()
    w, h = i.box(5, 5)

    prs = Presentation()
    prs.print_layout()

    prs.add_slide('title_and_body', title='hi', text=i)
    prs.save('/tmp/xdat/test.pptx')
