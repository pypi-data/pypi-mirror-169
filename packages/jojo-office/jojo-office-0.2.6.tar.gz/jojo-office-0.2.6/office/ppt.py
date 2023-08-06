#
# PPT : auto creation of the .pptx file
#
# Example:
# ::
#     import office
#
#     # create file output.pptx using template1.pptx
#     ppt = office.open_file("output.pptx", template="template1.pptx")
#
#     # create ppt content by fill the data from Excel file datafile.xlsx, and save.
#     ppt.fill('datafile.xlsx').save()
#
#     # save pptx to pdf with watermark (works on Windows with Microsoft PowerPoint)
#     ppt.save('final.pdf', watermark="CONFIDENTIAL")
#
#     # save pptx slides into a long image, with watermark (works on Windows with Microsoft PowerPoint)
#     ppt.save('long.jpg', watermark="CONFIDENTIAL")
#
#     # play ppt (works on Windows with Microsoft PowerPoint)
#     ppt.play()
#

import os
import io
import time
import copy
import tempfile
from .base import BaseFile
from .util import Util
from .variable import Variable, VarData
from .win32 import Win32
from .pdf import PDF
from .photo import Photo

try:
    import win32api
except ImportError:
    win32api = None

try:
    import playsound
except ImportError:
    playsound = None

import openpyxl
import pptx
from pptx.slide import Slides
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.shapes.picture import Picture
from pptx.shapes.graphfrm import GraphicFrame
from pptx.chart.chart import Chart
from pptx.shapes.shapetree import Shape
from pptx.chart.data import ChartData


class PPTConsts:
    # https://docs.microsoft.com/en-us/office/vba/api/powerpoint.ppsaveasfiletype
    ppSaveAsAddIn = 8
    ppSaveAsAnimatedGIF = 40
    ppSaveAsBMP = 19
    ppSaveAsDefault = 11
    ppSaveAsEMF = 23
    ppSaveAsExternalConverter = 64000
    ppSaveAsGIF = 16
    ppSaveAsJPG = 17
    ppSaveAsMetaFile = 15
    ppSaveAsMP4 = 39
    ppSaveAsOpenDocumentPresentation = 35
    ppSaveAsOpenXMLAddin = 30
    ppSaveAsOpenXMLPicturePresentation = 36
    ppSaveAsOpenXMLPresentation = 24
    ppSaveAsOpenXMLPresentationMacroEnabled = 25
    ppSaveAsOpenXMLShow = 28
    ppSaveAsOpenXMLShowMacroEnabled = 29
    ppSaveAsOpenXMLTemplate = 26
    ppSaveAsOpenXMLTemplateMacroEnabled = 27
    ppSaveAsOpenXMLTheme = 31
    ppSaveAsPDF = 32
    ppSaveAsPNG = 18
    ppSaveAsPresentation = 1
    ppSaveAsRTF = 6
    ppSaveAsShow = 7
    ppSaveAsStrictOpenXMLPresentation = 38
    ppSaveAsTemplate = 5
    ppSaveAsTIF = 21
    ppSaveAsWMV = 37
    ppSaveAsXMLPresentation = 34
    ppSaveAsXPS = 33


# PPT .pptx Document
class PPT(BaseFile):

    def _new_file(self, filename, **kwargs):
        self._obj = pptx.Presentation()
        self.filename = None
        if filename:
            self.save(filename, **kwargs)

    def _save_file(self, filename, **kwargs):
        formats = {
            '.pdf': PPTConsts.ppSaveAsPDF,
            '.ppt': PPTConsts.ppSaveAsPresentation,
            '.rtf': PPTConsts.ppSaveAsRTF,
            '.pot': PPTConsts.ppSaveAsTemplate,
            '.potx': PPTConsts.ppSaveAsTemplate,
        }

        file_ext = Util.file_ext(filename)

        if file_ext == '.pptx':
            return self._save_pptx(filename, **kwargs)

        elif file_ext in ['.jpg', '.jpeg']:
            return self._save_one_image(filename, **kwargs)

        elif file_ext in formats:
            app, document = Win32.open_powerpoint(self.filename, False)
            if not Util.is_absolute_path(filename):
                filename = os.path.join(os.getcwd(), filename)

            password = kwargs.pop('password', None)
            watermark = kwargs.pop('watermark', None)

            document.SaveAs(filename, formats[file_ext], **kwargs)
            Win32.close(app, document)

            if (password or watermark) and Util.file_ext(filename) == '.pdf':
                PDF(filename).watermark(watermark).save(password=password)

        else:
            return super()._save_file(filename, **kwargs)

    def _open(self, filename, **kwargs):
        self._obj = pptx.Presentation(filename)
        self.filename = filename

    def _close(self, **kwargs):
        self._obj = None
        return True

    def _save_pptx(self, filename, **kwargs):
        self._obj.save(filename, **kwargs)
        self.filename = filename

    @property
    def slide_height(self) -> int:
        return self._obj.slide_height

    @slide_height.setter
    def slide_height(self, height):
        self._obj.slide_height = height

    @property
    def slide_layouts(self):
        return self._obj.slide_layouts

    @property
    def slide_width(self) -> int:
        return self._obj.slide_width

    @slide_width.setter
    def slide_width(self, width):
        self._obj.slide_width = width

    def __len__(self):
        return len(self.slides)

    @property
    def slides(self) -> Slides:
        """ return a collection of slides """
        return self._obj.slides

    def _fill_data(self, data):
        """
        fill the data

        :param data:  dictionary object or Workbook object
        :return: self
        """

        i = 0  # 当前幻灯片索引号

        # 逐张幻灯片(Slide)处理
        while i < len(self._obj.slides):
            # 取得一张幻灯片
            slide = self._obj.slides[i]
            i += 1
            # 填充一张幻灯片
            ret = self._fill_slide(slide, data)
            if Variable.is_var(ret):
                # 处理复制幻灯片
                i = self._repeat_slide(i - 1, ret, data)
        return self

    def _repeat_slide(self, index, var_name, data: VarData):
        """ create repeat slides """
        if not data.has_var(var_name):
            return index + 1

        # 找到变量的最大行数
        repeat_count = data.var_length(var_name)
        if repeat_count == 0:
            self.delete_slide(index)
        elif repeat_count == 1:
            self._fill_slide(self.slides[index], data, 1)
            index += 1
        else:
            # 复制slide多次
            for row in range(repeat_count-1):
                self.clone_slide(index, index + row + 1)

            # 填充slide多次
            for row in range(0, repeat_count):
                self._fill_slide(self.slides[index], data, row + 1)
                index += 1

        return index

    def _fill_slide(self, slide, data, row=0):
        """ fill the data into the slide """
        # 逐个处理Shape
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:  # 如果是 TextFrame
                    # 对于每一格段落
                    for paragraph in shape.text_frame.paragraphs:
                        ret = self._fill_paragraph(slide, shape, paragraph, data, row)
                        if ret is not None:
                            return ret
                elif shape.has_table:  # 如果是 Table
                    self._fill_table(shape.table, data)
                elif isinstance(shape, pptx.shapes.picture.Picture):  # 如果是 Picture
                    pass
                elif shape.has_chart:  # 如果是 Chart
                    self._fill_chart(shape.chart, data)
                else:
                    pass
            except Exception as e:
                raise e

    def _find_paths(self, data):
        """ create path list for find files """
        paths = []
        if data:
            if isinstance(data, str):
                paths.append(data)
            elif isinstance(data, list):
                paths += data
            elif isinstance(data, tuple):
                paths += list(data)
            elif isinstance(data, VarData):
                paths.append(data.path)

        path = Util.file_path(self.filename)
        if path:
            paths.append(path)

        return paths

    def _fill_paragraph(self, slide, shape, paragraph, data: VarData, row=0):
        """ fill the data into the paragraph """

        def process(expr):
            """ 处理一个变量 """
            value = data.get_var(expr, row)
            value = str(value) if value is not None else ''
            if value.startswith(Variable.OBJECT_PREFIX):  # 如果 value 是对象
                # noinspection PyBroadException
                try:
                    filename = value[len(Variable.OBJECT_PREFIX):]
                    filename = Util.find_file(filename, self._find_paths(data))
                    # 插入对象
                    self._insert_object(slide, shape, filename)
                except Exception:
                    pass
                return ''
            else:
                return value

        var_name = ''
        for run in paragraph.runs:
            old_text = run.text
            text = ''
            for idx, c in enumerate(old_text):
                if len(var_name) == 0:
                    if c == '{':
                        var_name = c
                    else:
                        text += c
                else:
                    if c == '}':
                        var_name += c
                        v = Variable(var_name)
                        # 如果有 repeat 标志，且 repeat 未开始
                        if v.is_repeat and row == 0:
                            remain_text = old_text[idx + 1:]
                            # 将 varname 写入 run.text
                            run.text = text + var_name + remain_text
                            return var_name
                        else:
                            if data.has_var(v):
                                s = process(v)
                                text += s
                            else:
                                text += var_name
                        var_name = ''
                    elif c == '{':
                        text += var_name
                        var_name = c
                    else:
                        var_name += c
            run.text = text

    # noinspection PyMethodMayBeStatic
    def _insert_object(self, slide, shape, filename):
        """ insert object into ppt， object is load from file, object's position is same as the shape"""
        if Util.is_image(filename):
            slide.shapes.add_picture(filename, shape.left, shape.top, shape.width, shape.height)
        elif Util.is_video(filename):
            poster = None
            slide.shapes.add_movie(filename, shape.left, shape.top, shape.width, shape.height,
                                   poster_frame_image=poster)
        elif Util.is_audio(filename):
            pass  # TODO

    # noinspection PyMethodMayBeStatic
    def _fill_table(self, table, data: VarData):
        """ fill the data into table """
        def fill_text_frame(frame, text):
            filled = False
            for paragraph in frame.paragraphs:
                if len(paragraph.runs) == 0:
                    if not filled:
                        run = paragraph.add_run()
                        run.text = text
                        if paragraph.font.size:
                            run.font.size = paragraph.font.size
                        filled = True
                else:
                    for run in paragraph.runs:
                        if not filled:
                            run.text = text
                            filled = True
                        else:
                            run.text = ''
            return filled

        if len(table.rows) == 0:
            return

        # 逐列处理
        for col in range(len(table.columns)):
            # 取得该列第一行的文字
            word = str(table.rows[0].cells[col].text).strip()
            # 如果是变量
            if Variable.is_var(word):
                # 逐行处理, 纵向填充数据
                for row in range(len(table.rows)):
                    value = data.get_var(word, row) if data.has_var(word) else word
                    value = str(value) if value is not None else ''
                    cell = table.rows[row].cells[col]
                    if hasattr(cell, 'text_frame'):
                        if not fill_text_frame(cell.text_frame, value):
                            cell.text = value
                    else:
                        cell.text = value

    def _fill_chart(self, chart: Chart, data):
        """ fill the data into the chart """
        table = self._get_chart_values(chart, data)
        if len(table) > 0 and len(table[0]) > 0:
            chart_data = ChartData()
            chart_data.categories = table[0][1:]
            for i in range(1, len(table)):
                chart_data.add_series(str(table[i][0]), table[i][1:])
            chart.replace_data(chart_data)

    # noinspection PyMethodMayBeStatic
    def _get_chart_headers(self, chart: Chart):
        """
        get a list of field from the first row of the datasheet of the chart.

        :Chinese: 取得 chart 数据表的第一行的值的列表， 返回 list

        :param chart:
        :return:
        """
        result = []
        # noinspection PyBroadException
        try:
            wb = openpyxl.load_workbook(io.BytesIO(chart.part.chart_workbook.xlsx_part.blob))
            if len(wb.sheetnames) >= 1:
                sheet_name = wb.sheetnames[0]
                sheet = wb[sheet_name]
                row = col = 1
                while True:
                    val = sheet.cell(row, col).value
                    if val is None or val == '':
                        break
                    else:
                        result.append(str(val))
                        col += 1
        except Exception:
            pass
        return result

    # noinspection PyMethodMayBeStatic
    def _set_chart_headers(self, chart: Chart, headers):
        """
        set the chart headers.

        :Chinese: 设置 chart 数据表的第一行的值.

        :param chart:
        :param headers:
        :return: None
        """
        wb = openpyxl.load_workbook(io.BytesIO(chart.part.chart_workbook.xlsx_part.blob))
        if len(wb.sheetnames) >= 1:
            sheet_name = wb.sheetnames[0]
            sheet = wb[sheet_name]
            row = col = 1
            # set header values
            for val in headers:
                sheet.cell(row, col).value = val
                col += 1
            # update chart workbook
            buffer = io.BytesIO()
            wb.save(buffer)
            chart.part.chart_workbook.update_from_xlsx_blob(buffer.getvalue())

    def _get_chart_values(self, chart: Chart, data: VarData):
        """
        get the values for the chart.

        :Chinese: 取得 chart 所需的数据

        :param chart:
        :param data:
        :return:
        """
        headers = self._get_chart_headers(chart)
        chart_values = []
        for index, word in enumerate(headers):
            word = str(word).strip()
            if Variable.is_var(word):
                values = []
                row = 0
                while True:
                    val = data.get_var(word, row) if data.has_var(word) else word
                    if val is None or val == '' or Variable.is_var(val):
                        break
                    else:
                        values.append(val)
                        row += 1
                chart_values.append(values)
        return chart_values

    def _clone_chart(self, shape, dst_slide, left=None, top=None,
                     width=None, height=None):
        """ clone the chart """
        if not shape.has_chart:
            return None

        left = shape.left if left is None else left
        top = shape.top if top is None else top
        width = shape.width if width is None else width
        height = shape.height if height is None else height

        # create chart data from shape.chart
        chart_data = ChartData()
        plot = None
        for p in shape.chart.plots:
            plot = p
            chart_data.categories = p.categories
            for series in p.series:
                chart_data.add_series(series.name, series.values)
            break

        # add chart
        # noinspection PyProtectedMember
        shapes = shape._parent if dst_slide is None else dst_slide.shapes
        new_shape = shapes.add_chart(
            shape.chart.chart_type, left, top, width, height, chart_data
        )

        headers = self._get_chart_headers(shape.chart)
        self._set_chart_headers(new_shape.chart, headers)

        if hasattr(plot, 'gap_width'):
            new_shape.chart.plots[0].gap_width = plot.gap_width
            new_shape.chart.plots[0].overlap = plot.overlap

        if hasattr(plot, 'bubble_scale'):
            new_shape.chart.plots[0].bubble_scale = plot.bubble_scale

        if plot:
            new_shape.chart.plots[0].has_data_labels = plot.has_data_labels
            new_shape.chart.plots[0].vary_by_categories = plot.vary_by_categories

        return new_shape

    def clone_shape(self, shape, dst_slide, left=None, top=None,
                    width=None, height=None):
        """ colon the shape """
        left = shape.left if left is None else left
        top = shape.top if top is None else top
        width = shape.width if width is None else width
        height = shape.height if height is None else height

        if shape.has_chart:
            return self._clone_chart(shape, dst_slide, left, top, width, height)

        elif isinstance(shape, pptx.shapes.picture.Picture):
            # noinspection PyProtectedMember
            shapes = shape._parent if dst_slide is None else dst_slide.shapes
            return shapes.add_picture(io.BytesIO(shape.image.blob), left, top, width, height)

        else:
            # https://github.com/scanny/python-pptx/issues/533
            if dst_slide is None:
                sp_tree = shape.element.getparent()
                parent = None
            else:
                sp_tree = dst_slide.shapes.element
                parent = dst_slide.shapes

            new_sp = copy.deepcopy(shape.element)
            sp_tree.append(new_sp)

            if isinstance(shape, GraphicFrame):
                new_shape = GraphicFrame(new_sp, parent)
            else:
                new_shape = Shape(new_sp, parent)
            # new_shape.element._nvXxPr.cNvPr.id += 1000
            # new_shape.element._nvXxPr.cNvPr.name = "SomeThing"
            new_shape.left = left
            new_shape.top = top
            new_shape.width = width
            new_shape.height = height

            return new_shape

    # noinspection PyMethodMayBeStatic
    def delete_shape(self, shape):
        """ delete shape """
        if shape:
            shape.element.getparent().remove(shape.element)

    @property
    def _xml_slides(self):
        """ return slides xml """
        # https://github.com/scanny/python-pptx/issues/274
        # noinspection PyProtectedMember
        return self._obj.slides._sldIdLst

    def move_slide(self, old_index, new_index):
        """ move the slide to a new position. """
        slides = list(self._xml_slides)
        self._xml_slides.remove(slides[old_index])
        self._xml_slides.insert(new_index, slides[old_index])

    def delete_slide(self, index):
        """ delete slide at specified index """
        slides = list(self._xml_slides)
        self._xml_slides.remove(slides[index])

    def clone_slide(self, index, new_index=None):
        """ clone the slide """
        src_slide = self._obj.slides[index]
        # create new slide
        dst_slide = self._obj.slides.add_slide(src_slide.slide_layout)

        # empty the dst_slide, delete all shapes in dst_slide
        count = 0
        max_count = len(dst_slide.shapes)
        while len(dst_slide.shapes) > 0 and count < max_count:
            self.delete_shape(dst_slide.shapes[0])
            count += 1

        # copy shapes from src_slide
        for shape in src_slide.shapes:
            self.clone_shape(shape, dst_slide)

        # move dst_slide to new position
        if new_index is not None:
            old_index = self._obj.slides.index(dst_slide)
            self.move_slide(old_index, new_index)

        return dst_slide

    def slides_to_images(self, path):
        """
        convert all slide to images.

        :param path: the file path(dirctory) where the image files to save.
        :return: path
        """
        app, presentation = Win32.powerpoint_open(self.filename)

        if path.find(':') <= 0:
            path = os.path.join(os.getcwd(), path)

        presentation.SaveAs(path, 17)  # save each slide to image

        Win32.powerpoint_close(app, presentation)
        return path

    # noinspection PyMethodMayBeStatic
    def _load_play_file(self, filename):
        """ load play parameters from file. """
        result = []
        f = open(filename, 'r')
        lines = f.readlines()
        f.close()

        for idx, line in enumerate(lines):
            line = line.strip()
            if line:
                c = line[:1]
                if c == '#':  # 注释行
                    continue
                if '0' < c < '9':
                    offset = line.find(' ')
                    if offset < 0:
                        offset = line.find('\t')
                    if offset < 0:
                        result.append(int(line))
                    else:
                        num = line[:offset].strip()
                        file = line[offset+1:].strip()
                        if file.rfind('#') >= 0:
                            file = file[:file.rfind('#')].strip()
                        if len(file) > 0:
                            result.append([int(num), file])
                        else:
                            result.append(int(num))
                else:
                    raise ValueError('file %s format error at line %s' % (repr(filename), idx+1))
        return result

    def play(self, interval=4, loop=True):
        """
        play the Powerpoint.

        :param interval:  (optional) interval time in seconds. could be float or list of float or filename
            :Chinese: (可选)翻页的时间间隔(秒)，可以是 float 或 float数组 或 文件名
        :param loop:     (optional) whether loop
        :return:
        """
        def get_data(i, offset=0):
            """ 取得offset位置的数据， offset=0 为翻页的时间间隔(秒) """
            data = interval
            default_interval = 4
            if isinstance(data, int) or isinstance(data, float):
                return data if offset == 0 else None
            elif isinstance(data, list):
                if 0 <= i < len(data):
                    item = data[i]
                    if isinstance(item, list) or isinstance(item, tuple):
                        if len(item) > offset:
                            val = item[offset]
                            if offset == 0:
                                if not (isinstance(val, int) or isinstance(val, float)):
                                    return default_interval
                            return val
                return default_interval if offset == 0 else None

        def sleep(i):
            """  休眠  """
            seconds = get_data(i)  # 取得 翻页的时间间隔(秒)
            # time.sleep(seconds)
            win32api.Sleep(int(seconds * 1000))

        def sound(i, path=''):
            """ 播放声音 """
            filename = get_data(i, 1)
            if Util.is_audio(filename):
                ext = Util.file_ext(filename)
                filename = Util.find_file(filename, self._find_paths(path))
                Util.check_module(playsound, 'playsound')
                # noinspection PyBroadException
                try:
                    playsound.playsound(filename)
                except Exception:
                    pass

        play_file_path = ''
        if isinstance(interval, str):
            interval = self._load_play_file(interval)
            play_file_path = Util.file_path(interval)

        # open powerpoint app
        app, presentation = Win32.powerpoint_open(self.filename)

        # https://devblogs.microsoft.com/scripting/how-can-i-run-a-powerpoint-slide-show-from-a-script
        slide_count = presentation.Slides.Count
        run = presentation.SlideShowSettings.Run()
        view = presentation.SlideShowWindow.View
        # presentation.SlideShowWindow.View.Next()

        Win32.bring_to_front(app)
        canceled = False
        index = 0
        while not canceled:
            # noinspection PyBroadException
            try:
                sound(index, play_file_path)  # 播放当前页的声音
                sleep(index)   # 当前页的休眠

                slide_index = view.Slide.SlideIndex
                if slide_index < slide_count:
                    index += 1
                    view.Next()
                elif loop:
                    index = 0
                    view.First()
                else:
                    canceled = True
            except Exception:
                # raise(e)
                canceled = True

        Win32.powerpoint_close(app, presentation)

    # noinspection PyMethodMayBeStatic
    def _get_slide_image_filelist(self, path):
        filenames = []  # 图片文件名列表
        file_prefix = ''  # 图片文件前缀
        file_ext = ''  # 图片文件扩展名

        # 分析图片文件， 识别文件前缀、文件扩展名
        for fn in os.listdir(path):
            ext = Util.file_ext(fn, False)
            if ext.lower() in ['.jpg', '.jpeg']:
                file_ext = ext
                filenames.append(os.path.join(path, fn))
                if file_prefix == '':
                    # 识别文件前缀
                    fname = fn[:len(fn) - len(ext)]
                    for c in fname:
                        if '0' <= c <= '9':
                            break
                        else:
                            file_prefix += c

        # 将图片文件排序
        img_files = []
        for n in range(1, len(filenames) + 1):
            fn = os.path.join(path, file_prefix + str(n) + file_ext)
            if fn in filenames:
                img_files.append(os.path.join(path, fn))
            else:
                print('Missing slide', str(n))

        return img_files

    def _save_one_image(self, jpg_filename, **kwargs):
        """
        Save all slides to a long image.

        :Chinese: 将PPT保存为一张长图(jpg文件)

        :param jpg_filename:  JPG filename of the long image.
        :return: return True if success.
        """
        # 创建临时目录
        temp_dir = tempfile.TemporaryDirectory().name
        # 将幻灯片存为图片
        temp_dir = self.slides_to_images(temp_dir)
        # 读取图片文件列表
        filenames = self._get_slide_image_filelist(temp_dir)

        # 读取图片文件, 拼接为长图
        if len(filenames) > 0:
            watermark = kwargs.get('watermark', None)
            result = None
            try:
                height = 0
                for j, fn in enumerate(filenames):
                    img = Photo(fn)
                    if watermark:
                        img.watermark(watermark)
                    if j == 0:
                        width, height = img.size  # 获取拼接图片的宽和高
                        result = Photo((width, height * len(filenames)), mode=img.mode)
                    result.paste(img, box=(0, j * height))
                if result:
                    result.save(jpg_filename)
                    return True
            finally:
                Util.delete_dir(temp_dir)
        return False

    def print(self, **kwargs):
        """ print the file to default printer. """
        print("printing", self.filename, '...')

        app, presentation = Win32.open_powerpoint(self._absolute_filename, False)
        presentation.PrintOut(**kwargs)
        time.sleep(2)
        Win32.close(app, presentation)
